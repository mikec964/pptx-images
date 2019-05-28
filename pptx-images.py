#!/usr/bin/env python3

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def save_pic(picture, suffix, out_dir=''):
    '''Save picture to disk'''

    image = picture.image
    image_bytes = image.blob
    image_filename = 'image ' + suffix + '.' + image.ext
    print(f'saving {suffix}')
    with open(os.path.join(out_dir,image_filename), 'wb') as f:
        f.write(image_bytes)


def pics_in_shapes(shape_list):
    '''Return pics in list of shapes

    Shapes can be picture, group, text_box, other.

    '''
    pic_list = []
    for shape in shape_list.shapes:
        # add pics to the list of pics
        # add groups to the list of groups
        # call the group-walker
        if shape == MSO_SHAPE_TYPE.PICTURE:
            pic_list.append(shape)
        elif shape == MSO_SHAPE_TYPE.GROUP:
            pic_list.extend(pics_in_shapes(shape))
    return pic_list


def pics_in_pres(filename, out_dir):
    '''Extract pics from pptx presentation to disk'''

    pres = Presentation(filename)
    nSlides = len(pres.slides)
    slideN = 0
    for slide in pres.slides:
        slideN += 1
        nShapes = len(slide.shapes)
        shapeN = 0
        for shape in slide.shapes:
            shapeN += 1
            print(f'slide {slideN} of {nSlides}, shape {shapeN} of {nShapes} is {shape.shape_type}')
            suffix = str(slideN) + '-' + str(shapeN)
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                save_pic(shape, suffix, out_dir)
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for shp in shape.shapes:
                    print(shp.shape_type)
                    if shp.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        save_pic(shp, suffix, out_dir)



def main():
    filename = 'ARG032-Flat.pptx'
    out_dir = 'arg032-flat'
    pics_in_pres(filename, out_dir)
    return


if __name__ == "__main__":
    main()

