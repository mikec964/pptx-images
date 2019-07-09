import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


class Pptimgx:
    '''Extracts images from a PPTX

    pres1 = Pptimgx('ARG049-Horizontal.pptx')
    print(f'file: {pres1.filename}, slides: {pres1.nSlides}, pics: {pres1.nPics}')
    for p in range(pres1.nPics):
        pres1.save(p, path='arg049')   # 'image 1.png'
    n = 1
    with open(f'image {n}.{pres1.pics[n].image.ext}', 'wb') as f:
        f.write(pres1.pics[n].image.blob)
    '''

    def __init__(self, filename):
        self.filename = filename
        self.pres = Presentation(filename)
        self.nSlides = len(self.pres.slides)
        self.pics = self.extract_pictures(self.pres)
        self.nPics = len(self.pics)


    def __repr__(self):
        return f'Presentation {self.filename} has {self.nSlides} slides, {len(self.pics)} pics.'


    def extract_pictures(self, pres):
        pics = []
        slideN = 0
        for slide in self.pres.slides:
            slideN += 1
            nShapes = len(slide.shapes)
            shapeN = 0
            for shape in slide.shapes:
                shapeN += 1
                # This should be recursive but only goes 1 layer deep
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    pics.append(shape)
                elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    for shape in shape.shapes:
                        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                            pics.append(shape)
        return pics


    def save(self, pic_num, prefix='image ', suffix='', path=''):
        image = self.pics[pic_num].image
        image_bytes = image.blob
        image_filename = prefix + str(pic_num) + suffix + '.' + image.ext
        print(f'saving {image_filename}')
        with open(os.path.join(path,image_filename), 'wb') as f:
            f.write(image_bytes)

